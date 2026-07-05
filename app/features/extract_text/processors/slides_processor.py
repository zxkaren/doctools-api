from pathlib import Path

from pptx import Presentation


def extract_text_from_table(table) -> str:
    extracted_table_lines = []

    for table_row in table.rows:
        table_cells_text = []

        for table_cell in table_row.cells:
            cell_text = table_cell.text.strip()

            if cell_text:
                table_cells_text.append(cell_text)

        if table_cells_text:
            extracted_table_lines.append(" | ".join(table_cells_text))

    return "\n".join(extracted_table_lines)


def extract_text_from_shape(shape) -> str:
    extracted_shape_blocks = []

    if getattr(shape, "has_text_frame", False):
        shape_text = shape.text.strip()

        if shape_text:
            extracted_shape_blocks.append(shape_text)

    if getattr(shape, "has_table", False):
        table_text = extract_text_from_table(shape.table)

        if table_text:
            extracted_shape_blocks.append(table_text)

    # Recursividade para capturar textos dentro de grupos de shapes.
    if hasattr(shape, "shapes"):
        for nested_shape in shape.shapes:
            nested_shape_text = extract_text_from_shape(nested_shape)

            if nested_shape_text:
                extracted_shape_blocks.append(nested_shape_text)

    return "\n".join(extracted_shape_blocks)


def extract_text_from_slides(file_path: Path) -> str:
    """
    Resumo:
        Extrai apenas textos visíveis de um arquivo PPTX, incluindo caixas de
        texto, placeholders e tabelas, sem capturar imagens, gráficos ou notas.

    Parâmetros:
        file_path (Path): caminho do arquivo PPTX salvo em storage.

    Retorno:
        str: texto bruto extraído da apresentação.
    """
    extracted_slide_blocks = []

    try:
        presentation = Presentation(file_path)

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text_blocks = []

            for shape in slide.shapes:
                shape_text = extract_text_from_shape(shape)

                if shape_text:
                    slide_text_blocks.append(shape_text)

            if slide_text_blocks:
                extracted_slide_blocks.append(
                    f"Slide {slide_number}\n" + "\n".join(slide_text_blocks)
                )

        return "\n\n".join(extracted_slide_blocks)

    except Exception as error:
        print(f"Erro ao extrair texto do PPTX {file_path.name}: {error}")
        raise