import os
from io import BytesIO
from pathlib import Path

import fitz
import pytest
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches


os.environ["FLASK_ENV"] = "testing"
os.environ["DEBUG"] = "False"
os.environ["HOST"] = "127.0.0.1"
os.environ["PORT"] = "5000"
os.environ["TIMEZONE_NAME"] = "America/Sao_Paulo"
os.environ["SECRET_KEY"] = "test-secret-key"
os.environ["MAX_CONTENT_LENGTH"] = "52428800"

os.environ["STORAGE_ROOT"] = "storage"
os.environ["COMPARE_STORAGE_FOLDER"] = "compare"
os.environ["COMPARE_RECEIVED_FOLDER"] = "received"
os.environ["COMPARE_PROCESSED_FOLDER"] = "processed"
os.environ["COMPARE_TEMP_FOLDER"] = "temp"

os.environ["ALLOWED_COMPARE_EXTENSIONS"] = "pdf,docx,xlsx,pptx"
os.environ["IMPLEMENTED_COMPARE_EXTENSIONS"] = "pdf,docx,xlsx,pptx"

os.environ["DEFAULT_COMPARE_RESPONSE_MODE"] = "json_file"
os.environ["ALLOWED_COMPARE_RESPONSE_MODES"] = "download_url,json,json_file"

os.environ["CLEANUP_FILE_MAX_AGE_HOURS"] = "24"
os.environ["CLEANUP_INTERVAL_MINUTES"] = "60"


from app import create_app
from app.config import Config


@pytest.fixture
def client(tmp_path: Path):
    Config.COMPARE_RECEIVED_FOLDER = tmp_path / "storage" / "compare" / "received"
    Config.COMPARE_PROCESSED_FOLDER = tmp_path / "storage" / "compare" / "processed"
    Config.COMPARE_TEMP_FOLDER = tmp_path / "storage" / "compare" / "temp"

    Config.ALLOWED_COMPARE_EXTENSIONS = ["pdf", "docx", "xlsx", "pptx"]
    Config.IMPLEMENTED_COMPARE_EXTENSIONS = ["pdf", "docx", "xlsx", "pptx"]

    flask_app = create_app()
    flask_app.config["TESTING"] = True

    return flask_app.test_client()


def build_pdf_bytes(text: str) -> bytes:
    """
    Resumo:
        Cria um PDF em memória para uso nos testes da API.

    Parâmetros:
        text (str): texto que será inserido no PDF.

    Retorno:
        bytes: conteúdo binário do PDF gerado.
    """
    document = fitz.open()
    page = document.new_page()
    page.insert_text(fitz.Point(72, 72), text, fontsize=12)

    pdf_bytes = document.tobytes()
    document.close()

    return pdf_bytes


def build_xlsx_bytes(
    cell_values: dict[str, object],
    sheet_name: str = "Sheet1",
) -> bytes:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = sheet_name

    for cell_reference, cell_value in cell_values.items():
        worksheet[cell_reference] = cell_value

    workbook_bytes = BytesIO()
    workbook.save(workbook_bytes)

    return workbook_bytes.getvalue()


def add_pptx_text_slide(
    presentation: Presentation,
    text: str,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    text_box = slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(6),
        Inches(1),
    )
    text_box.text_frame.text = text


def build_pptx_bytes(
    first_slide_text: str,
    second_slide_text: str | None = None,
) -> bytes:
    presentation = Presentation()

    add_pptx_text_slide(
        presentation,
        first_slide_text,
    )

    if second_slide_text:
        add_pptx_text_slide(
            presentation,
            second_slide_text,
        )

    presentation_bytes = BytesIO()
    presentation.save(presentation_bytes)

    return presentation_bytes.getvalue()


def get_slide_text(slide) -> str:
    text_values = []

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text_values.append(shape.text_frame.text)

    return "\n".join(text_values)


def get_slide_notes(slide) -> str:
    return slide.notes_slide.notes_text_frame.text or ""


def get_presentation_notes(presentation: Presentation) -> str:
    return "\n".join(
        get_slide_notes(slide)
        for slide in presentation.slides
    )


def test_compare_pdf_returns_download_url_and_summary_table(client) -> None:
    original_pdf = build_pdf_bytes("hello world")
    modified_pdf = build_pdf_bytes("hello brave world")

    response = client.post(
        "/compare/",
        data={
            "original": (BytesIO(original_pdf), "original.pdf"),
            "modified": (BytesIO(modified_pdf), "modified.pdf"),
            "response_mode": "json_file",
        },
        content_type="multipart/form-data",
    )

    response_payload = response.get_json()

    assert response.status_code == 200
    assert response_payload["success"] is True
    assert "download_url" in response_payload["data"]
    assert "summary_table" in response_payload["data"]
    assert response_payload["data"]["summary_table"]["add"] >= 1
    assert response_payload["data"]["summary_table"]["total_changes"] >= 1


def test_compare_pdf_without_modified_file_returns_bad_request(client) -> None:
    original_pdf = build_pdf_bytes("hello world")

    response = client.post(
        "/compare/",
        data={
            "original": (BytesIO(original_pdf), "original.pdf"),
            "response_mode": "json_file",
        },
        content_type="multipart/form-data",
    )

    response_payload = response.get_json()

    assert response.status_code == 400
    assert response_payload["success"] is False
    assert response_payload["error"]["message"] == "informe o arquivo modified"


def test_compare_with_different_extensions_returns_bad_request(client) -> None:
    original_pdf = build_pdf_bytes("hello world")

    response = client.post(
        "/compare/",
        data={
            "original": (BytesIO(original_pdf), "original.pdf"),
            "modified": (BytesIO(b"fake content"), "modified.docx"),
            "response_mode": "json_file",
        },
        content_type="multipart/form-data",
    )

    response_payload = response.get_json()

    assert response.status_code == 400
    assert response_payload["success"] is False
    assert response_payload["error"]["message"] == "os arquivos devem ter a mesma extensão"


def test_compare_specific_route_with_wrong_extension_returns_bad_request(client) -> None:
    original_pdf = build_pdf_bytes("hello world")
    modified_pdf = build_pdf_bytes("hello brave world")

    response = client.post(
        "/compare/docx",
        data={
            "original": (BytesIO(original_pdf), "original.pdf"),
            "modified": (BytesIO(modified_pdf), "modified.pdf"),
            "response_mode": "json_file",
        },
        content_type="multipart/form-data",
    )

    response_payload = response.get_json()

    assert response.status_code == 400
    assert response_payload["success"] is False
    assert response_payload["error"]["message"] == "essa rota só aceita a extensão selecionada"


def test_compare_xlsx_returns_download_url_summary_and_processed_file(client) -> None:
    original_xlsx = build_xlsx_bytes(
        {
            "A1": "Item",
            "B1": "Value",
            "B2": "=SUM(B3:B4)",
            "B3": 10,
            "B4": 20,
            "C1": "Remove me",
        }
    )
    modified_xlsx = build_xlsx_bytes(
        {
            "A1": "Item",
            "B1": "Value",
            "B2": "=SUM(B3:B5)",
            "B3": 10,
            "B4": 20,
            "B5": 30,
        }
    )

    response = client.post(
        "/compare/",
        data={
            "original": (BytesIO(original_xlsx), "original.xlsx"),
            "modified": (BytesIO(modified_xlsx), "modified.xlsx"),
            "response_mode": "json_file",
        },
        content_type="multipart/form-data",
    )

    response_payload = response.get_json()

    assert response.status_code == 200
    assert response_payload["success"] is True
    assert "download_url" in response_payload["data"]
    assert "summary_table" in response_payload["data"]

    summary_table = response_payload["data"]["summary_table"]

    assert summary_table["add"] >= 1
    assert summary_table["delete"] >= 1
    assert summary_table["total_changes"] == summary_table["add"] + summary_table["delete"]

    processed_files = list(Config.COMPARE_PROCESSED_FOLDER.glob("*compared*.xlsx"))

    assert len(processed_files) == 1

    compared_workbook = load_workbook(processed_files[0], data_only=False)
    compared_worksheet = compared_workbook["Sheet1"]

    assert "summary_table" in compared_workbook.sheetnames
    assert compared_workbook["summary_table"]["A1"].value == "add"
    assert compared_workbook["summary_table"]["A2"].value == "delete"
    assert compared_workbook["summary_table"]["A3"].value == "total_changes"
    assert compared_worksheet["B5"].fill.fgColor.rgb.endswith("ADD8E6")
    assert compared_worksheet["B5"].comment is not None
    assert compared_worksheet["B5"].comment.text.startswith("add:")
    assert compared_worksheet["C1"].fill.fgColor.rgb.endswith("FF9393")
    assert compared_worksheet["C1"].comment is not None
    assert compared_worksheet["C1"].comment.text.startswith("delete:")


def test_compare_pptx_returns_download_url_summary_and_processed_file(client) -> None:
    original_pptx = build_pptx_bytes(
        first_slide_text="hello legacy world",
    )
    modified_pptx = build_pptx_bytes(
        first_slide_text="hello brave world",
        second_slide_text="new slide",
    )

    response = client.post(
        "/compare/",
        data={
            "original": (BytesIO(original_pptx), "original.pptx"),
            "modified": (BytesIO(modified_pptx), "modified.pptx"),
            "response_mode": "json_file",
        },
        content_type="multipart/form-data",
    )

    response_payload = response.get_json()

    assert response.status_code == 200
    assert response_payload["success"] is True
    assert "download_url" in response_payload["data"]
    assert "summary_table" in response_payload["data"]

    summary_table = response_payload["data"]["summary_table"]

    assert summary_table["add"] >= 2
    assert summary_table["delete"] >= 1
    assert summary_table["total_changes"] == summary_table["add"] + summary_table["delete"]

    processed_files = list(Config.COMPARE_PROCESSED_FOLDER.glob("*compared*.pptx"))

    assert len(processed_files) == 1

    compared_presentation = Presentation(processed_files[0])
    first_slide_text = get_slide_text(compared_presentation.slides[0])
    all_notes = get_presentation_notes(compared_presentation)
    last_slide_text = get_slide_text(compared_presentation.slides[-1])

    assert "hello brave world" in first_slide_text
    assert "hello legacy world" not in first_slide_text

    assert "DocTools Compare" in all_notes
    assert "delete: legacy" in all_notes
    assert "add: brave" in all_notes
    assert "add: slide adicionado" in all_notes

    assert "summary table" in last_slide_text