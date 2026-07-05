from io import BytesIO
from pathlib import Path

import fitz
import pytest
from docx import Document
from pptx import Presentation

from app import create_app
from app.config import Config
from app.features.extract_text.processors.text_cleaner import clean_extracted_text

@pytest.fixture
def client():
    flask_app = create_app()
    flask_app.config.update(TESTING=True)

    with flask_app.test_client() as test_client:
        yield test_client

def create_pdf_bytes(text_content: str) -> BytesIO:
    pdf_document = fitz.open()
    pdf_page = pdf_document.new_page()
    pdf_page.insert_text((72, 72), text_content)

    pdf_bytes = pdf_document.write()
    pdf_document.close()

    return BytesIO(pdf_bytes)


def create_docx_bytes(text_content: str) -> BytesIO:
    word_document = Document()
    word_document.add_paragraph(text_content)

    memory_file = BytesIO()
    word_document.save(memory_file)
    memory_file.seek(0)

    return memory_file


def create_pptx_bytes(text_content: str) -> BytesIO:
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(blank_layout)

    text_box = slide.shapes.add_textbox(100, 100, 500, 120)
    text_box.text = text_content

    memory_file = BytesIO()
    presentation.save(memory_file)
    memory_file.seek(0)

    return memory_file


def configure_extract_text_storage(monkeypatch, temporary_path: Path) -> None:
    received_folder = temporary_path / "received"
    processed_folder = temporary_path / "processed"
    temp_folder = temporary_path / "temp"

    received_folder.mkdir()
    processed_folder.mkdir()
    temp_folder.mkdir()

    monkeypatch.setattr(Config, "EXTRACT_TEXT_RECEIVED_FOLDER", received_folder)
    monkeypatch.setattr(Config, "EXTRACT_TEXT_PROCESSED_FOLDER", processed_folder)
    monkeypatch.setattr(Config, "EXTRACT_TEXT_TEMP_FOLDER", temp_folder)


def test_clean_extracted_text_removes_urls_emails_emojis_page_numbers_and_image_captions():
    raw_text_content = """
    Página 1
    Este é um texto válido para leitura.
    contato@empresa.com
    https://empresa.com.br
    Texto com emoji 😊 e pontuação preservada!
    Figura 1: foto da capa do livro
    Imagem 2 - Fluxo do processo
    Figure 3: Architecture diagram
    2/10
    """

    cleaned_text_content = clean_extracted_text(raw_text_content)

    assert "Página 1" not in cleaned_text_content
    assert "contato@empresa.com" not in cleaned_text_content
    assert "https://empresa.com.br" not in cleaned_text_content
    assert "😊" not in cleaned_text_content
    assert "Figura 1: foto da capa do livro" not in cleaned_text_content
    assert "Imagem 2 - Fluxo do processo" not in cleaned_text_content
    assert "Figure 3: Architecture diagram" not in cleaned_text_content
    assert "Texto com emoji e pontuação preservada!" in cleaned_text_content

def test_extract_text_returns_one_output_for_each_uploaded_file(client, monkeypatch, tmp_path):
    configure_extract_text_storage(monkeypatch, tmp_path)

    response = client.post(
        "/extract-text/",
        data={
            "output_format": "txt",
            "files": [
                (create_pdf_bytes("Texto extraído do PDF."), "documento.pdf"),
                (create_docx_bytes("Texto extraído do DOCX."), "documento.docx"),
                (create_pptx_bytes("Texto extraído do PPTX."), "documento.pptx"),
            ],
        },
        content_type="multipart/form-data",
    )

    response_json = response.get_json()
    response_data = response_json["data"]

    assert response.status_code == 200
    assert response_data["total_files"] == 3
    assert response_data["processed_files"] == 3
    assert response_data["failed_files"] == 0
    assert len(response_data["files"]) == 3

    for file_result in response_data["files"]:
        output_filename = file_result["output_filename"]
        output_path = Config.EXTRACT_TEXT_PROCESSED_FOLDER / output_filename

        assert file_result["status"] == "success"
        assert output_path.exists()
        assert output_path.suffix == ".txt"


def test_extract_text_uses_docx_as_default_output(client, monkeypatch, tmp_path):
    configure_extract_text_storage(monkeypatch, tmp_path)

    response = client.post(
        "/extract-text/",
        data={
            "files": [
                (create_docx_bytes("Texto usando saída padrão."), "padrao.docx"),
            ],
        },
        content_type="multipart/form-data",
    )

    response_json = response.get_json()
    response_data = response_json["data"]
    file_result = response_data["files"][0]
    output_path = Config.EXTRACT_TEXT_PROCESSED_FOLDER / file_result["output_filename"]

    assert response.status_code == 200
    assert response_data["output_format"] == "docx"
    assert output_path.exists()
    assert output_path.suffix == ".docx"


def test_extract_text_rejects_invalid_extension(client):
    response = client.post(
        "/extract-text/",
        data={
            "output_format": "txt",
            "files": [
                (BytesIO(b"invalid content"), "planilha.xlsx"),
            ],
        },
        content_type="multipart/form-data",
    )

    assert response.status_code == 400


def test_extract_text_rejects_invalid_output_format(client):
    response = client.post(
        "/extract-text/",
        data={
            "output_format": "pdf",
            "files": [
                (create_docx_bytes("Texto válido."), "documento.docx"),
            ],
        },
        content_type="multipart/form-data",
    )

    assert response.status_code == 400