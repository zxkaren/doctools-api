from collections import Counter
from difflib import SequenceMatcher
from hashlib import sha256
from pathlib import Path
from typing import Any
import re

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


ADD_TYPE = "add"
DELETE_TYPE = "delete"
TEXT_KIND = "text"
PICTURE_KIND = "picture"

NOTE_HEADER = "DocTools Compare"
SUMMARY_TITLE = "summary table"

MATCH_THRESHOLD = 0.45
POSITION_BONUS = 0.03
MAX_NOTE_TEXT_LENGTH = 500

SUMMARY_TITLE_FONT_SIZE = Pt(24)
SUMMARY_TABLE_FONT_SIZE = Pt(14)

IMAGE_ADDED_NOTE = "add: imagem adicionada"
IMAGE_REMOVED_NOTE = "delete: imagem removida"
SLIDE_ADDED_NOTE = "add: slide adicionado"
SLIDE_REMOVED_NOTE = "delete: slide removido"


def collect_shapes_recursively(shapes: Any) -> list[Any]:
    collected_shapes = []

    for shape in shapes:
        collected_shapes.append(shape)

        if hasattr(shape, "shapes"):
            collected_shapes.extend(collect_shapes_recursively(shape.shapes))

    return collected_shapes


def get_shapes_by_kind(slide: Any | None, shape_kind: str) -> list[Any]:
    if slide is None:
        return []

    slide_shapes = collect_shapes_recursively(slide.shapes)

    if shape_kind == TEXT_KIND:
        return [
            shape
            for shape in slide_shapes
            if getattr(shape, "has_text_frame", False)
        ]

    return [
        shape
        for shape in slide_shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]


def get_blank_layout(presentation: Any) -> Any:
    for slide_layout in presentation.slide_layouts:
        if slide_layout.name.lower() == "blank":
            return slide_layout

    return presentation.slide_layouts[-1]


def get_slide_text(slide: Any | None) -> str:
    if slide is None:
        return ""

    text_shapes = sorted(
        get_shapes_by_kind(slide, TEXT_KIND),
        key=lambda shape: (shape.top, shape.left),
    )

    text_values = []

    for shape in text_shapes:
        shape_text = (shape.text_frame.text or "").strip()

        if shape_text:
            text_values.append(shape_text)

    return "\n".join(text_values)


def get_words(text: str) -> list[str]:
    return re.findall(r"\S+", text or "")


def get_text_tokens(text: str) -> list[str]:
    return re.findall(r"\S+\s*|\s+", text or "")


def normalize_text_for_similarity(text: str) -> str:
    normalized_text = " ".join(get_words(text.lower()))

    return normalized_text.strip()


def clean_changed_text(text: str) -> str:
    cleaned_text = " ".join(get_words(text))

    if len(cleaned_text) <= MAX_NOTE_TEXT_LENGTH:
        return cleaned_text

    return f"{cleaned_text[:MAX_NOTE_TEXT_LENGTH]}..."


def get_picture_hashes(slide: Any | None) -> Counter:
    return Counter(
        sha256(picture_shape.image.blob).hexdigest()
        for picture_shape in get_shapes_by_kind(slide, PICTURE_KIND)
    )


def build_slide_signature(slide: Any | None) -> str:
    text_signature = normalize_text_for_similarity(get_slide_text(slide))
    picture_signatures = " ".join(
        sorted(get_picture_hashes(slide).elements())
    )

    return f"{text_signature} {picture_signatures}".strip()


def build_slide_records(presentation: Any) -> list[dict[str, Any]]:
    slide_records = []

    for slide_index, slide in enumerate(presentation.slides):
        slide_records.append(
            {
                "index": slide_index,
                "slide": slide,
                "signature": build_slide_signature(slide),
            }
        )

    return slide_records


def calculate_slide_score(
    original_record: dict[str, Any],
    modified_record: dict[str, Any],
) -> float:
    original_signature = original_record["signature"]
    modified_signature = modified_record["signature"]

    if not original_signature and not modified_signature:
        return 1.0

    similarity_score = SequenceMatcher(
        None,
        original_signature,
        modified_signature,
        autojunk=False,
    ).ratio()

    if original_record["index"] == modified_record["index"]:
        similarity_score += POSITION_BONUS

    return min(similarity_score, 1.0)


def match_slides_by_similarity(
    original_records: list[dict[str, Any]],
    modified_records: list[dict[str, Any]],
) -> dict[int, int]:
    """
    Resumo:
        Pareia slides originais e modificados por similaridade de conteúdo.

    Parâmetros:
        original_records (list[dict[str, Any]]): slides do arquivo original.
        modified_records (list[dict[str, Any]]): slides do arquivo modificado.

    Retorno:
        dict[int, int]: mapa com índice do slide original e índice do slide modificado.
    """
    candidate_matches = []

    for original_record in original_records:
        for modified_record in modified_records:
            similarity_score = calculate_slide_score(
                original_record,
                modified_record,
            )

            if similarity_score >= MATCH_THRESHOLD:
                candidate_matches.append(
                    (
                        similarity_score,
                        original_record["index"],
                        modified_record["index"],
                    )
                )

    candidate_matches.sort(reverse=True)
    matched_original_indexes = set()
    matched_modified_indexes = set()
    slide_matches = {}

    for similarity_score, original_index, modified_index in candidate_matches:
        if original_index in matched_original_indexes:
            continue

        if modified_index in matched_modified_indexes:
            continue

        slide_matches[original_index] = modified_index
        matched_original_indexes.add(original_index)
        matched_modified_indexes.add(modified_index)

    return slide_matches


def build_text_note_lines(
    original_text: str,
    modified_text: str,
) -> tuple[list[str], int, int]:
    original_tokens = get_text_tokens(original_text)
    modified_tokens = get_text_tokens(modified_text)

    sequence_matcher = SequenceMatcher(
        None,
        original_tokens,
        modified_tokens,
        autojunk=False,
    )

    note_lines = []
    add_count = 0
    delete_count = 0

    for operation, original_start, original_end, modified_start, modified_end in sequence_matcher.get_opcodes():
        original_fragment = "".join(original_tokens[original_start:original_end])
        modified_fragment = "".join(modified_tokens[modified_start:modified_end])

        if operation in ("delete", "replace"):
            deleted_text = clean_changed_text(original_fragment)

            if deleted_text:
                note_lines.append(f"{DELETE_TYPE}: {deleted_text}")
                delete_count += 1

        if operation in ("insert", "replace"):
            added_text = clean_changed_text(modified_fragment)

            if added_text:
                note_lines.append(f"{ADD_TYPE}: {added_text}")
                add_count += 1

    return note_lines, add_count, delete_count


def build_image_note_lines(
    original_slide: Any | None,
    modified_slide: Any | None,
) -> tuple[list[str], int, int]:
    original_hashes = get_picture_hashes(original_slide)
    modified_hashes = get_picture_hashes(modified_slide)

    note_lines = []
    add_count = 0
    delete_count = 0

    for picture_hash, modified_quantity in modified_hashes.items():
        original_quantity = original_hashes[picture_hash]
        added_quantity = max(modified_quantity - original_quantity, 0)

        if added_quantity:
            note_lines.extend([IMAGE_ADDED_NOTE] * added_quantity)
            add_count += added_quantity

    for picture_hash, original_quantity in original_hashes.items():
        modified_quantity = modified_hashes[picture_hash]
        deleted_quantity = max(original_quantity - modified_quantity, 0)

        if deleted_quantity:
            note_lines.extend([IMAGE_REMOVED_NOTE] * deleted_quantity)
            delete_count += deleted_quantity

    return note_lines, add_count, delete_count


def write_slide_notes(slide: Any, note_lines: list[str]) -> None:
    if not note_lines:
        return

    notes_text_frame = slide.notes_slide.notes_text_frame
    current_notes = (notes_text_frame.text or "").strip()
    compare_notes = "\n".join([NOTE_HEADER, *note_lines])

    notes_text_frame.text = (
        f"{current_notes}\n\n{compare_notes}"
        if current_notes
        else compare_notes
    )


def compare_matched_slides(
    original_slide: Any,
    modified_slide: Any,
) -> tuple[int, int]:
    text_note_lines, text_add_count, text_delete_count = build_text_note_lines(
        get_slide_text(original_slide),
        get_slide_text(modified_slide),
    )
    image_note_lines, image_add_count, image_delete_count = build_image_note_lines(
        original_slide,
        modified_slide,
    )

    note_lines = [*text_note_lines, *image_note_lines]
    write_slide_notes(modified_slide, note_lines)

    return (
        text_add_count + image_add_count,
        text_delete_count + image_delete_count,
    )


def compare_presentations(
    original_presentation: Any,
    modified_presentation: Any,
) -> tuple[int, int, list[str]]:
    """
    Resumo:
        Compara apresentações preservando visualmente os slides existentes.

    Parâmetros:
        original_presentation (Any): apresentação original aberta pelo python-pptx.
        modified_presentation (Any): apresentação modificada aberta pelo python-pptx.

    Retorno:
        tuple[int, int, list[str]]: totais de add, delete e notas globais.
    """
    original_records = build_slide_records(original_presentation)
    modified_records = build_slide_records(modified_presentation)
    slide_matches = match_slides_by_similarity(
        original_records,
        modified_records,
    )

    add_count = 0
    delete_count = 0
    summary_notes = []

    matched_original_indexes = set(slide_matches.keys())
    matched_modified_indexes = set(slide_matches.values())

    for original_index, modified_index in sorted(slide_matches.items()):
        slide_add_count, slide_delete_count = compare_matched_slides(
            original_records[original_index]["slide"],
            modified_records[modified_index]["slide"],
        )

        add_count += slide_add_count
        delete_count += slide_delete_count

    for modified_record in modified_records:
        if modified_record["index"] in matched_modified_indexes:
            continue

        write_slide_notes(
            modified_record["slide"],
            [f"{SLIDE_ADDED_NOTE}: {modified_record['index'] + 1}"],
        )
        add_count += 1

    for original_record in original_records:
        if original_record["index"] in matched_original_indexes:
            continue

        summary_notes.append(
            f"{SLIDE_REMOVED_NOTE}: {original_record['index'] + 1}"
        )
        delete_count += 1

    return add_count, delete_count, summary_notes


def add_summary_slide(
    presentation: Any,
    summary_table: dict[str, int],
    summary_notes: list[str],
) -> None:
    summary_slide = presentation.slides.add_slide(get_blank_layout(presentation))

    title_box = summary_slide.shapes.add_textbox(
        Inches(1),
        Inches(0.6),
        Inches(8),
        Inches(0.6),
    )
    title_run = title_box.text_frame.paragraphs[0].add_run()
    title_run.text = SUMMARY_TITLE
    title_run.font.bold = True
    title_run.font.size = SUMMARY_TITLE_FONT_SIZE

    table = summary_slide.shapes.add_table(
        4,
        2,
        Inches(1),
        Inches(1.6),
        Inches(7),
        Inches(2),
    ).table

    summary_rows = [
        ("change_type", "count"),
        ("add", str(summary_table["add"])),
        ("delete", str(summary_table["delete"])),
        ("total_changes", str(summary_table["total_changes"])),
    ]

    for row_index, summary_row in enumerate(summary_rows):
        change_label, change_count = summary_row

        for column_index, cell_text in enumerate((change_label, change_count)):
            text_frame = table.cell(row_index, column_index).text_frame
            text_frame.clear()

            text_run = text_frame.paragraphs[0].add_run()
            text_run.text = cell_text
            text_run.font.size = SUMMARY_TABLE_FONT_SIZE

    write_slide_notes(summary_slide, summary_notes)


def compare_files(
    original_path: Path,
    modified_path: Path,
    processed_path: Path,
) -> dict[str, dict[str, int]]:
    """
    Resumo:
        Compara dois arquivos PPTX preservando visualmente os slides modificados.

    Parâmetros:
        original_path (Path): caminho do arquivo PPTX original.
        modified_path (Path): caminho do arquivo PPTX modificado.
        processed_path (Path): caminho onde o arquivo comparado será salvo.

    Retorno:
        dict[str, dict[str, int]]: resumo com add, delete e total_changes.
    """
    try:
        if original_path.suffix.lower() != ".pptx" or modified_path.suffix.lower() != ".pptx":
            raise ValueError("apenas arquivos .pptx são suportados pelo motor de slides")

        original_presentation = Presentation(original_path)
        modified_presentation = Presentation(modified_path)

        add_count, delete_count, summary_notes = compare_presentations(
            original_presentation,
            modified_presentation,
        )

        summary_table = {
            "add": add_count,
            "delete": delete_count,
            "total_changes": add_count + delete_count,
        }

        add_summary_slide(
            modified_presentation,
            summary_table,
            summary_notes,
        )

        processed_path.parent.mkdir(parents=True, exist_ok=True)
        modified_presentation.save(processed_path)

        print("comparação pptx concluída com sucesso")

        return {
            "summary_table": summary_table,
        }

    except Exception as error:
        print(f"erro ao comparar pptx: {error}")
        raise